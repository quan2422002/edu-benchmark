#!/usr/bin/env python3
"""Build the VietInfoTutorBench KSE 2026 paper-finalization meeting deck."""
from __future__ import annotations

import math
import tempfile
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "kse_submit_manuscript" / "slide"
FIG_DIR = ROOT / "kse_submit_manuscript" / "manuscript" / "figures"
OUT_PPTX = OUT_DIR / "VietInfoTutorBench_KSE_2026_paper_finalization.pptx"

# 16:9 canvas
W, H = 13.333, 7.5
FONT = "Noto Sans"
FONT_MONO = "DejaVu Sans Mono"

NAVY = "102A43"
BLUE = "2563EB"
TEAL = "0F766E"
CYAN = "0891B2"
GREEN = "16A34A"
AMBER = "D97706"
RED = "DC2626"
PURPLE = "7C3AED"
INK = "17202A"
MUTED = "52606D"
LIGHT = "F5F7FA"
LIGHT_BLUE = "EAF2FF"
LIGHT_TEAL = "E8F5F2"
LIGHT_AMBER = "FFF4E5"
LIGHT_RED = "FDECEC"
WHITE = "FFFFFF"
BORDER = "D9E2EC"
DARK_BG = "0B1F33"

SECTION_COLORS = {
    "Câu chuyện": BLUE,
    "Xây dựng benchmark": TEAL,
    "Khung đánh giá": PURPLE,
    "Thực nghiệm": AMBER,
    "Kết luận": RED,
    "Phụ lục": MUTED,
}


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_cell_border(cell, color=BORDER, width="12700"):
    # python-pptx has no public table-border API.
    from lxml import etree
    from pptx.oxml.ns import qn
    tcPr = cell._tc.get_or_add_tcPr()
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        tag = tcPr.find(qn(edge))
        if tag is None:
            tag = etree.SubElement(tcPr, qn(edge))
        tag.set("w", width)
        solid = tag.find(qn("a:solidFill"))
        if solid is None:
            solid = etree.SubElement(tag, qn("a:solidFill"))
        srgb = solid.find(qn("a:srgbClr"))
        if srgb is None:
            srgb = etree.SubElement(solid, qn("a:srgbClr"))
        srgb.set("val", color)


def add_rect(slide, x, y, w, h, fill, line=None, radius=False, transparency=0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    if transparency:
        shp.fill.transparency = transparency
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(0.8)
    else:
        shp.line.fill.background()
    return shp


def add_line(slide, x1, y1, x2, y2, color=BORDER, width=1.5, dash=None):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = dash
    return line


def set_text_style(run, size=18, color=INK, bold=False, font=FONT, italic=False):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=FONT,
             margin=0.05, italic=False, line_spacing=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    set_text_style(run, size, color, bold, font, italic)
    return box


def add_rich_lines(slide, lines, x, y, w, h, size=18, color=INK,
                   bullet=False, bullet_color=None, spacing=7, margin=0.04):
    """lines: [(text, bold, color)] or strings; one paragraph per item."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    for idx, item in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        p.line_spacing = 1.05
        if bullet:
            p.text = "• "
            set_text_style(p.runs[0], size, bullet_color or BLUE, True)
        parts = item if isinstance(item, list) else [(item, False, color)]
        for text, bold, part_color in parts:
            r = p.add_run()
            r.text = text
            set_text_style(r, size, part_color or color, bold)
    return box


def add_title(slide, title, section="Câu chuyện", subtitle=None, dark=False):
    if dark:
        add_rect(slide, 0, 0, W, H, DARK_BG)
        color = WHITE
        add_text(slide, title, 0.75, 0.78, 11.9, 1.15, 30, color, True, valign=MSO_ANCHOR.MIDDLE)
        if subtitle:
            add_text(slide, subtitle, 0.78, 1.92, 11.7, 0.8, 17, "D9E2EC")
        return
    add_rect(slide, 0, 0, W, 0.16, SECTION_COLORS.get(section, BLUE))
    add_text(slide, section.upper(), 0.7, 0.28, 3.4, 0.35, 10, SECTION_COLORS.get(section, BLUE), True)
    add_text(slide, title, 0.7, 0.66, 11.95, 0.72, 28, NAVY, True, valign=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, subtitle, 0.72, 1.34, 11.8, 0.45, 13, MUTED)


def add_footer(slide, number, source=None, dark=False):
    color = "9FB3C8" if dark else "7B8794"
    add_line(slide, 0.7, 7.14, 12.63, 7.14, "35506B" if dark else BORDER, 0.7)
    add_text(slide, "VietInfoTutorBench • KSE 2026", 0.72, 7.18, 4.0, 0.22, 8.5, color)
    if source:
        add_text(slide, source, 4.15, 7.18, 7.75, 0.22, 7.8, color, align=PP_ALIGN.RIGHT)
    add_text(slide, str(number), 12.15, 7.17, 0.45, 0.24, 8.5, color, True, align=PP_ALIGN.RIGHT)


def add_takeaway(slide, text, color=BLUE, y=6.54):
    add_rect(slide, 0.7, y, 11.93, 0.46, color, radius=True)
    add_text(slide, "THÔNG ĐIỆP CHỐT", 0.9, y+0.08, 1.65, 0.22, 9, WHITE, True, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, text, 2.45, y+0.055, 9.95, 0.30, 12, WHITE, True, valign=MSO_ANCHOR.MIDDLE)


def add_card(slide, x, y, w, h, title, body=None, accent=BLUE, fill=WHITE,
             title_size=15, body_size=12.5, number=None):
    add_rect(slide, x, y, w, h, fill, BORDER, radius=True)
    add_rect(slide, x, y, 0.08, h, accent, radius=True)
    if number is not None:
        add_rect(slide, x+0.22, y+0.20, 0.48, 0.48, accent, radius=True)
        add_text(slide, str(number), x+0.22, y+0.20, 0.48, 0.48, 14, WHITE, True,
                 PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        tx = x+0.82
        tw = w-1.02
    else:
        tx = x+0.26
        tw = w-0.45
    add_text(slide, title, tx, y+0.18, tw, 0.42, title_size, NAVY, True)
    if body:
        if isinstance(body, list):
            add_rich_lines(slide, body, x+0.28, y+0.72, w-0.52, h-0.87,
                           body_size, MUTED, bullet=True, spacing=5)
        else:
            add_text(slide, body, x+0.28, y+0.70, w-0.52, h-0.84, body_size, MUTED)


def add_stat(slide, x, y, w, h, value, label, accent=BLUE, note=None):
    add_rect(slide, x, y, w, h, WHITE, BORDER, radius=True)
    add_text(slide, value, x+0.14, y+0.14, w-0.28, 0.55, 26, accent, True, PP_ALIGN.CENTER)
    add_text(slide, label, x+0.16, y+0.70, w-0.32, 0.42, 12.5, NAVY, True, PP_ALIGN.CENTER)
    if note:
        add_text(slide, note, x+0.14, y+1.12, w-0.28, h-1.21, 9.5, MUTED, False, PP_ALIGN.CENTER)


def add_picture_contain(slide, path, x, y, w, h, pad=0.0):
    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min((w-2*pad)/iw, (h-2*pad)/ih)
    pw, ph = iw*scale, ih*scale
    px = x + (w-pw)/2
    py = y + (h-ph)/2
    slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def add_arrow(slide, x, y, w=0.55, h=0.34, color=BLUE):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(color)
    shp.line.fill.background()
    return shp


def add_table(slide, x, y, w, h, data, col_widths=None, header_fill=NAVY,
              font_size=10.5, first_col_bold=False, row_fills=None):
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(w * cw / total)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.05)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            fill = header_fill if r == 0 else (row_fills[r-1] if row_fills and r-1 < len(row_fills) else WHITE)
            cell.fill.solid(); cell.fill.fore_color.rgb = rgb(fill)
            cell.text = str(data[r][c])
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            p.vertical_anchor = MSO_ANCHOR.MIDDLE
            for run in p.runs:
                set_text_style(run, font_size if r else font_size, WHITE if r == 0 else INK,
                               r == 0 or (first_col_bold and c == 0))
            set_cell_border(cell)
    return table


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid(); bg.fill.fore_color.rgb = rgb(LIGHT)
    return slide


def setup_chart_style():
    plt.rcParams.update({
        "font.family": "DejaVu Sans",
        "font.size": 10,
        "axes.titleweight": "bold",
        "axes.edgecolor": "#CBD5E1",
        "axes.labelcolor": "#334155",
        "xtick.color": "#475569",
        "ytick.color": "#475569",
    })


def save_overall_chart(path):
    setup_chart_style()
    labels = ["Gemini\nbaseline", "Gemini +\nLearnLM", "Llama 4\nMaverick"]
    gem = [87.87, 85.35, 49.51]
    gpt = [84.01, 84.78, 74.56]
    x = np.arange(3); width = 0.34
    fig, ax = plt.subplots(figsize=(8.8, 4.5), dpi=180)
    b1 = ax.bar(x-width/2, gem, width, label="Gemini judge", color="#2563EB")
    b2 = ax.bar(x+width/2, gpt, width, label="GPT judge", color="#D97706")
    ax.set_ylim(0, 105); ax.set_ylabel("Overall Accuracy (%)")
    ax.set_xticks(x, labels); ax.grid(axis="y", alpha=.2); ax.set_axisbelow(True)
    ax.legend(frameon=False, ncol=2, loc="upper center")
    for bars in (b1,b2):
        ax.bar_label(bars, fmt="%.2f", padding=3, fontsize=9, fontweight="bold")
    for s in ("top","right"): ax.spines[s].set_visible(False)
    fig.tight_layout(); fig.savefig(path, transparent=True, bbox_inches="tight"); plt.close(fig)


def save_forest_chart(path):
    setup_chart_style()
    names = [
        "Gemini: Baseline − LearnLM", "Gemini: Baseline − Llama", "Gemini: LearnLM − Llama",
        "GPT: Baseline − LearnLM", "GPT: Baseline − Llama", "GPT: LearnLM − Llama",
    ]
    vals = np.array([2.52, 38.36, 35.84, -0.77, 9.45, 10.22])
    lows = np.array([1.53,35.53,32.89,-2.14,7.43,8.04])
    highs = np.array([3.57,41.21,38.82,0.45,11.71,13.04])
    colors = ["#2563EB"]*3+["#D97706"]*3
    y = np.arange(len(names))[::-1]
    fig, ax = plt.subplots(figsize=(9.2, 4.5), dpi=180)
    for yi,v,lo,hi,c in zip(y, vals, lows, highs, colors):
        ax.errorbar(v, yi, xerr=[[v-lo],[hi-v]], fmt="o", color=c, ecolor=c,
                    elinewidth=2.5, capsize=4, markersize=7)
        ax.text(hi+0.7, yi, f"{v:+.2f} [{lo:+.2f}, {hi:+.2f}]", va="center", fontsize=8.7)
    ax.axvline(0, color="#64748B", lw=1.2, ls="--")
    ax.set_yticks(y, names); ax.set_xlabel("Chênh lệch Overall Accuracy (điểm phần trăm)")
    ax.set_xlim(-5, 46); ax.grid(axis="x", alpha=.18); ax.set_axisbelow(True)
    for s in ("top","right","left"): ax.spines[s].set_visible(False)
    fig.tight_layout(); fig.savefig(path, transparent=True, bbox_inches="tight"); plt.close(fig)


def save_ablation_chart(path):
    setup_chart_style()
    cats = ["Holistic", "General", "Questioning", "Overall"]
    gem = np.array([-0.21,-3.45,-4.25,-2.52]); gem_lo=np.array([-1.65,-4.82,-6.69,-3.58]); gem_hi=np.array([1.20,-2.06,-1.77,-1.50])
    gpt = np.array([1.07,0.77,0.87,0.77]); gpt_lo=np.array([-0.44,-0.75,-1.67,-0.48]); gpt_hi=np.array([2.64,2.29,3.44,2.14])
    y=np.arange(len(cats))[::-1]
    fig,ax=plt.subplots(figsize=(8.6,4.1),dpi=180)
    for offset,vals,los,his,c,label in [(0.12,gem,gem_lo,gem_hi,"#2563EB","Gemini judge"),(-0.12,gpt,gpt_lo,gpt_hi,"#D97706","GPT judge")]:
        ax.errorbar(vals,y+offset,xerr=[vals-los,his-vals],fmt="o",color=c,ecolor=c,
                    capsize=3,elinewidth=2,markersize=6,label=label)
    ax.axvline(0,color="#64748B",lw=1.1,ls="--")
    ax.set_yticks(y,cats); ax.set_xlabel("Δ = LearnLM-oriented − baseline (điểm %)")
    ax.set_xlim(-8,5); ax.grid(axis="x",alpha=.18); ax.legend(frameon=False,ncol=2,loc="lower center")
    for s in ("top","right","left"): ax.spines[s].set_visible(False)
    fig.tight_layout(); fig.savefig(path,transparent=True,bbox_inches="tight"); plt.close(fig)


def save_agreement_chart(path):
    setup_chart_style()
    labels=["Tất cả", "Gemini baseline", "Gemini + LearnLM", "Llama 4 Maverick"]
    exact=[80.45,88.36,89.79,63.21]; ac1=[77.5,87.4,89.0,52.9]
    x=np.arange(4); width=.34
    fig,ax=plt.subplots(figsize=(8.6,4.2),dpi=180)
    b1=ax.bar(x-width/2,exact,width,label="Exact agreement",color="#0F766E")
    b2=ax.bar(x+width/2,ac1,width,label="Gwet's AC1 × 100",color="#7C3AED")
    ax.set_ylim(0,105); ax.set_ylabel("Điểm / tỷ lệ (%)"); ax.set_xticks(x,labels)
    ax.grid(axis="y",alpha=.18); ax.set_axisbelow(True); ax.legend(frameon=False,ncol=2,loc="upper center")
    for bars in (b1,b2): ax.bar_label(bars,fmt="%.1f",padding=2,fontsize=8.5)
    for s in ("top","right"): ax.spines[s].set_visible(False)
    fig.tight_layout(); fig.savefig(path,transparent=True,bbox_inches="tight"); plt.close(fig)


def save_criterion_chart(path):
    setup_chart_style()
    labels=["Giao tiếp chung","Chính xác chung","Questioning","Hỗ trợ chung","Bám trạng thái","Feedback","Practice","Explanation","Modelling","Challenge"]
    vals=[61.90,65.55,68.12,68.98,72.90,76.88,77.78,84.69,86.74,95.83]
    counts=[4200,4200,5859,4200,4200,7254,243,7767,837,72]
    colors=["#DC2626" if v<70 else "#D97706" if v<80 else "#0F766E" for v in vals]
    y=np.arange(len(labels))
    fig,ax=plt.subplots(figsize=(8.8,5.0),dpi=180)
    bars=ax.barh(y,vals,color=colors,height=.62)
    ax.set_xlim(50,102); ax.set_xlabel("Exact agreement (%)"); ax.set_yticks(y,labels)
    ax.grid(axis="x",alpha=.18); ax.set_axisbelow(True)
    for bar,v,n in zip(bars,vals,counts):
        ax.text(v+.5,bar.get_y()+bar.get_height()/2,f"{v:.2f}%  (N={n:,})",va="center",fontsize=8.3)
    for s in ("top","right","left"): ax.spines[s].set_visible(False)
    fig.tight_layout(); fig.savefig(path,transparent=True,bbox_inches="tight"); plt.close(fig)


def save_flags_chart(path):
    setup_chart_style()
    labels=["Feedback chỉ xác nhận/khen","Lập luận score 4–5 không nhất quán","Questioning không cần câu trả lời","Không có nguyên tắc ≥ 4","Tổ hợp thiếu hỗ trợ","Trên 3 nguyên tắc"]
    vals=[592,20,10,8,7,1]
    y=np.arange(len(labels))[::-1]
    fig,ax=plt.subplots(figsize=(8.7,3.8),dpi=180)
    bars=ax.barh(y,vals,color=["#D97706"]+["#94A3B8"]*5,height=.58)
    ax.set_yticks(y,labels); ax.set_xlabel("Số candidate bị cờ (các cờ có thể chồng lấn)")
    ax.grid(axis="x",alpha=.18); ax.set_axisbelow(True)
    for bar,v in zip(bars,vals): ax.text(v+5,bar.get_y()+bar.get_height()/2,str(v),va="center",fontweight="bold")
    ax.set_xlim(0,650)
    for s in ("top","right","left"): ax.spines[s].set_visible(False)
    fig.tight_layout(); fig.savefig(path,transparent=True,bbox_inches="tight"); plt.close(fig)


def build_deck():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    prs.core_properties.title = "VietInfoTutorBench — KSE 2026 paper-finalization meeting"
    prs.core_properties.subject = "Vietnamese lower-secondary Informatics AI tutor benchmark"
    prs.core_properties.author = "VietInfoTutorBench project team"
    prs.core_properties.comments = "Concise 12-slide deck generated from the latest KSE manuscript and experiment 20260727_170150."

    with tempfile.TemporaryDirectory(prefix="vietinfotutor-slides-") as td:
        td = Path(td)
        charts = {
            "overall": td/"overall.png",
            "forest": td/"forest.png",
            "agreement": td/"agreement.png",
        }
        save_overall_chart(charts["overall"])
        save_forest_chart(charts["forest"])
        save_agreement_chart(charts["agreement"])

        n = 0
        def finish(slide, source=None, dark=False):
            nonlocal n
            n += 1
            add_footer(slide, n, source, dark)
            return slide

        # 1. Title
        s = new_slide(prs); add_title(s, "VietInfoTutorBench", dark=True)
        add_text(s, "Benchmark gia sư AI Tin học THCS tiếng Việt\ntừ hội thoại do giáo viên biên soạn", 0.82, 2.0, 8.15, 1.45, 27, WHITE, True)
        add_rect(s, 0.85, 3.72, 5.85, 0.58, BLUE, radius=True)
        add_text(s, "THUYẾT TRÌNH CHỐT BÀI BÁO KSE 2026", 1.08, 3.86, 5.38, 0.25, 13, WHITE, True, PP_ALIGN.CENTER)
        add_text(s, "30/07/2026  •  Nhóm UET – HNMU", 0.88, 4.52, 6.6, 0.35, 15, "D9E2EC")
        x0=8.72
        for i,(v,l,c) in enumerate([("1.050","hội thoại",CYAN),("665","đạt audit",TEAL),("2.028","candidate",PURPLE),("1.400","eligible",AMBER)]):
            add_stat(s,x0+(i%2)*2.0,1.65+(i//2)*1.75,1.75,1.38,v,l,c)
        finish(s, dark=True)

        # 2. Problem, gaps, and contributions
        s=new_slide(prs); add_title(s,"Bài báo giải quyết vấn đề gì?","Câu chuyện")
        add_text(s,"BA KHOẢNG TRỐNG",.75,1.42,4.0,.34,14,BLUE,True)
        gaps=[
            ("Miền","Chủ yếu Toán/STEM rộng; thiếu Tin học THCS chuyên biệt.",BLUE),
            ("Ngôn ngữ","Đánh giá chủ yếu bằng tiếng Anh; ít hội thoại gia sư tiếng Việt.",TEAL),
            ("Dữ liệu","Ít quy trình bắt đầu từ hội thoại giáo viên thông thường và giữ truy vết.",AMBER),
        ]
        for i,(t,b,c) in enumerate(gaps):
            add_card(s,.72,1.86+i*1.20,5.45,.96,t,b,c,title_size=15,body_size=11.5,number=i+1)
        add_text(s,"BA ĐÓNG GÓP",6.62,1.42,4.0,.34,14,PURPLE,True)
        contributions=[
            ("Pipeline","Audit → candidate → requirement → deterministic lock.",TEAL),
            ("Đo lường","6 nguyên tắc + 6 năng lực → rubric 4 + 3n.",PURPLE),
            ("Benchmark & evidence","1.400 mẫu; 3 tutor × 2 judge.",AMBER),
        ]
        for i,(t,b,c) in enumerate(contributions):
            add_card(s,6.57,1.86+i*1.20,5.72,.96,t,b,c,title_size=15,body_size=11.5,number=i+1)
        add_rect(s,1.23,5.66,10.85,.48,LIGHT_BLUE,BORDER,radius=True)
        add_text(s,"Đánh giá gia sư là bài toán thiết kế phép đo sư phạm trong ngữ cảnh — không chỉ chấm đáp án đúng/sai.",1.45,5.78,10.4,.24,14,NAVY,True,PP_ALIGN.CENTER)
        add_takeaway(s,"Novelty nằm ở tổ hợp miền × ngôn ngữ × quy trình dữ liệu, không dùng tuyên bố “đầu tiên” tuyệt đối.",BLUE)
        finish(s,"Sections I–II")

        # 3. Overall pipeline
        s=new_slide(prs); add_title(s,"Pipeline xây dựng VietInfoTutorBench","Xây dựng benchmark")
        add_rect(s,.65,1.45,12.05,4.85,WHITE,BORDER,radius=True)
        add_picture_contain(s,FIG_DIR/"overall_pipeline_illustrated.png",.82,1.72,11.72,4.28)
        add_takeaway(s,"Phase 1 bảo đảm chất lượng dữ liệu; Phase 2 bảo đảm cấu trúc đo lường; Phase 3 kết hợp hai đầu vào thành benchmark.",TEAL)
        finish(s,"Figure 1 trong paper")

        # 4. Phase 1
        s=new_slide(prs); add_title(s,"Phase 1 — Kiểm toán 1.050 hội thoại giáo viên","Xây dựng benchmark")
        add_picture_contain(s,FIG_DIR/"phase_1_flow.png",.58,1.48,5.72,4.62)
        add_card(s,6.48,1.48,3.04,1.42,"Grounding học liệu",[
            "154 đơn vị SGK/SGV.","2.750 fragment + SQLite FTS."],TEAL,title_size=15,body_size=11.5)
        add_card(s,9.66,1.48,2.92,1.42,"Hybrid audit",[
            "Code + gpt-5.4-mini.","18 tiêu chí có evidence."],BLUE,title_size=15,body_size=11.5)
        add_stat(s,6.48,3.15,1.82,1.42,"665","pass",GREEN,"được conversion")
        add_stat(s,8.43,3.15,1.82,1.42,"382","review",AMBER,"không dùng hiện tại")
        add_stat(s,10.38,3.15,1.82,1.42,"3","failed",RED,"bị loại")
        add_rect(s,6.48,4.85,5.72,1.10,LIGHT_AMBER,BORDER,radius=True)
        add_text(s,"Strict aggregation",6.72,5.03,1.72,.28,14,AMBER,True)
        add_text(s,"Có fail → failed; không fail nhưng có uncertain → review; còn lại → pass.",8.25,4.98,3.7,.58,12.2,NAVY,True)
        add_takeaway(s,"Audit giữ truy vết tới học liệu và criterion; chỉ 665 hội thoại pass đi vào conversion.",GREEN)
        finish(s,"Section III.A")

        # 5. Phase 2
        s=new_slide(prs); add_title(s,"Phase 2 — Sáu nguyên tắc và sáu năng lực","Xây dựng benchmark")
        add_text(s,"NGUYÊN TẮC SƯ PHẠM",.85,1.43,5.35,.35,15,TEAL,True,PP_ALIGN.CENTER)
        principles=[("Challenge","Yêu cầu nhận thức"),("Explanation","Làm rõ ý/quan hệ"),("Modelling","Biểu diễn cách làm"),("Practice","Cơ hội luyện tập"),("Feedback","Nhận xét + cải thiện"),("Questioning","Hỏi để tiến triển")]
        for i,(a,b) in enumerate(principles):
            x=.78+(i%2)*2.75; y=1.92+(i//2)*1.04
            add_card(s,x,y,2.55,.82,a,b,TEAL,title_size=13.5,body_size=10.2)
        add_text(s,"NĂNG LỰC GIA SƯ",7.03,1.43,5.35,.35,15,PURPLE,True,PP_ALIGN.CENTER)
        caps=[("Accuracy","Đúng + có căn cứ"),("Learner state","Hiểu trạng thái HS"),("Strategy","Chọn chiến lược"),("Scaffolding","Điều tiết hỗ trợ"),("Diagnosis","Chẩn đoán lỗi"),("Communication","Rõ + đúng lứa tuổi")]
        for i,(a,b) in enumerate(caps):
            x=6.97+(i%2)*2.75; y=1.92+(i//2)*1.04
            add_card(s,x,y,2.55,.82,a,b,PURPLE,title_size=13.5,body_size=10.2)
        add_rect(s,5.98,2.14,.55,2.95,LIGHT_AMBER,BORDER,radius=True)
        add_text(s,"↔",6.02,3.23,.46,.46,24,AMBER,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
        add_rect(s,1.35,5.34,10.63,.63,WHITE,BORDER,radius=True)
        add_text(s,"Nguyên tắc kích hoạt yêu cầu theo mẫu; năng lực tổ chức coverage và observable evidence của rubric.",1.6,5.50,10.15,.28,14,NAVY,True,PP_ALIGN.CENTER)
        add_takeaway(s,"Ranh giới then chốt: state ≠ diagnosis; strategy ≠ scaffolding; explanation ≠ modelling.",PURPLE)
        finish(s,"Allison & Tharby; KMP-Bench; adaptive scaffolding; ECD")

        # 6. Phase 3
        s=new_slide(prs); add_title(s,"Phase 3 — Từ hội thoại thô đến 1.400 mẫu eligible","Xây dựng benchmark")
        add_card(s,.72,1.48,3.72,3.98,"1. Turn-level conversion",None,TEAL,title_size=17)
        bubbles=[("S1","Em không biết làm sao để ra 101...",LIGHT_BLUE,BLUE),("T1","Chia dãy 0–7; số 5 ở bên nào?",LIGHT_TEAL,TEAL),("S2","Bên phải; ghi bit 1.",LIGHT_BLUE,BLUE),("T2","Đúng! Chia tiếp dãy chứa số 5...",LIGHT_TEAL,TEAL)]
        for i,(r,t,fill,c) in enumerate(bubbles):
            y=2.18+i*.58
            add_rect(s,.96,y,3.22,.45,fill,BORDER,radius=True)
            add_text(s,r,1.08,y+.1,.35,.22,9.5,c,True)
            add_text(s,t,1.49,y+.08,2.5,.25,9.8,NAVY)
        add_text(s,"D1 → C1 (∅ history) + C2 (T1,S2) + C3 (...)",1.00,4.62,3.10,.45,11.5,PURPLE,True,PP_ALIGN.CENTER)
        add_card(s,4.80,1.48,3.52,3.98,"2. Requirement scoring",[
            "Input: context + source question + gold_answer.",
            "Không dùng gold_response.",
            "Chấm đủ 6 nguyên tắc theo thang 1–5.",
            "Code chọn score ≥ 4."],PURPLE,title_size=17,body_size=12.2)
        add_card(s,8.68,1.48,3.92,3.98,"3. Deterministic gates",[
            "1–3 nguyên tắc bắt buộc.",
            "Independent need + counterfactual.",
            "Không xung đột optionality.",
            "Tổ hợp đủ support theo candidate/family."],AMBER,title_size=17,body_size=12.2)
        add_rect(s,1.30,5.72,10.65,.48,LIGHT_RED,BORDER,radius=True)
        add_text(s,"2.028 candidate  →  12.168 scores  →  1.400 eligible / 655 family  +  628 bị loại",1.55,5.84,10.15,.25,14,NAVY,True,PP_ALIGN.CENTER)
        add_takeaway(s,"Model xử lý ngữ nghĩa; code xử lý threshold, validation, join và eligibility.",TEAL)
        finish(s,"Section III.C; candidate HNMU-G6-R0041-STT12")

        # 7. Pool distribution
        s=new_slide(prs); add_title(s,"Pool 1.400 mẫu: phân bố nguyên tắc và hội thoại","Xây dựng benchmark")
        add_rect(s,.62,1.43,12.1,4.85,WHITE,BORDER,radius=True)
        add_picture_contain(s,FIG_DIR/"phase3_candidate_statistics.png",.72,1.63,11.9,4.42)
        add_takeaway(s,"Pool phủ đủ lớp 6–9 và lịch sử đa lượt, nhưng Challenge (8) và Practice (27) còn rất hiếm.",AMBER)
        finish(s,"Figure 2 trong paper")

        # 8. Evaluation framework
        s=new_slide(prs); add_title(s,"Evaluation framework — từ context đến Win/Tie/Lose","Khung đánh giá")
        add_rect(s,.55,1.40,12.25,4.72,WHITE,BORDER,radius=True)
        add_picture_contain(s,FIG_DIR/"ai_tutor_prompting_and_judging.png",.68,1.58,11.98,4.30)
        chips=[("Native roles","USER ↔ ASSISTANT",BLUE),("Tutor không thấy","gold / rubric / ID",RED),("Rubric theo mẫu","4 chung + 3n riêng",PURPLE)]
        for i,(a,b,c) in enumerate(chips):
            x=.92+i*4.12
            add_rect(s,x,5.84,3.75,.36,WHITE,c,radius=True)
            add_text(s,a+": "+b,x+.1,5.90,3.55,.22,10.5,c,True,PP_ALIGN.CENTER)
        add_takeaway(s,"Judge so sánh mù tutor response với teacher reference theo từng rubric và một holistic judgment độc lập.",PURPLE)
        finish(s,"Section IV; 22 rubric đã được giáo viên HNMU xác nhận sơ bộ")

        # 9. Experiment setup
        s=new_slide(prs); add_title(s,"Thiết kế thực nghiệm","Thực nghiệm")
        add_text(s,"3 CẤU HÌNH AI TUTOR",.78,1.42,4.0,.34,15,AMBER,True)
        tutors=[("Gemini baseline","Gemini 3.5 Flash • baseline instruction",BLUE),("Gemini + LearnLM","Cùng base model • LearnLM-oriented instruction",TEAL),("Llama 4 Maverick","Llama MaaS • baseline instruction",PURPLE)]
        for i,(a,b,c) in enumerate(tutors): add_card(s,.76,1.88+i*1.10,4.42,.88,a,b,c,title_size=15,body_size=11.5)
        add_text(s,"2 LLM JUDGE",5.50,1.42,3.0,.34,15,AMBER,True)
        add_card(s,5.48,1.88,3.02,1.12,"Gemini 3.5 Flash","thinking level = medium",BLUE,title_size=15,body_size=12)
        add_card(s,5.48,3.22,3.02,1.12,"GPT-5.4-mini","reasoning effort = medium",AMBER,title_size=15,body_size=12)
        add_rect(s,5.48,4.57,3.02,1.15,LIGHT_AMBER,BORDER,radius=True)
        add_text(s,"Báo riêng từng judge",5.70,4.78,2.58,.28,15,NAVY,True,PP_ALIGN.CENTER)
        add_text(s,"Không lấy trung bình để che bất đồng",5.70,5.18,2.58,.26,10.5,MUTED,False,PP_ALIGN.CENTER)
        add_text(s,"QUY MÔ",8.92,1.42,3.0,.34,15,AMBER,True)
        add_stat(s,8.86,1.88,1.66,1.27,"1.400","samples",TEAL)
        add_stat(s,10.70,1.88,1.66,1.27,"4.200","responses",BLUE)
        add_stat(s,8.86,3.38,1.66,1.27,"8.400","overall",RED)
        add_stat(s,10.70,3.38,1.66,1.27,"77.664","rubric",PURPLE)
        add_rect(s,8.86,4.88,3.50,.82,LIGHT_BLUE,BORDER,radius=True)
        add_text(s,"5.000 cluster-bootstrap\nresamples theo sample_id",9.08,5.05,3.05,.42,12,NAVY,True,PP_ALIGN.CENTER)
        add_takeaway(s,"Ba target được so với cùng teacher reference trên cùng 1.400 sample và cùng rubric contract.",AMBER)
        finish(s,"Section V.A")

        # 10. Results
        s=new_slide(prs); add_title(s,"Kết quả chính — phân biệt lớn, chưa ổn định ở khác biệt nhỏ","Thực nghiệm")
        add_picture_contain(s,charts["overall"],.45,1.40,6.15,4.65)
        add_picture_contain(s,charts["forest"],6.65,1.38,6.15,4.72)
        add_rect(s,.95,5.82,3.64,.42,LIGHT_TEAL,GREEN,radius=True)
        add_text(s,"4/4 Gemini–Llama CI không chứa 0",1.08,5.91,3.38,.22,11.5,GREEN,True,PP_ALIGN.CENTER)
        add_rect(s,4.85,5.82,3.64,.42,LIGHT_AMBER,AMBER,radius=True)
        add_text(s,"Family-macro không đảo thứ hạng",4.98,5.91,3.38,.22,11.5,AMBER,True,PP_ALIGN.CENTER)
        add_rect(s,8.75,5.82,3.64,.42,LIGHT_RED,RED,radius=True)
        add_text(s,"Baseline–LearnLM đảo theo judge",8.88,5.91,3.38,.22,11.5,RED,True,PP_ALIGN.CENTER)
        add_takeaway(s,"Claim mạnh nhất: benchmark phân biệt bền vững khoảng cách lớn Gemini–Llama; chưa xếp hạng chắc hai Gemini configuration.",GREEN)
        finish(s,"Tables II–V; Overall Accuracy và 95% cluster-bootstrap CI")

        # 11. Judge agreement and threats
        s=new_slide(prs); add_title(s,"Judge agreement và threats to validity","Thực nghiệm")
        add_picture_contain(s,charts["agreement"],.48,1.42,6.18,4.52)
        add_card(s,6.88,1.48,2.72,1.18,"Bất đồng tập trung ở Llama","Exact agreement chỉ 63,21%.",RED,title_size=14.5,body_size=11.5)
        add_card(s,9.82,1.48,2.72,1.18,"Bất đối xứng","Gemini Lose/GPT Win: 419; chiều ngược: 65.",AMBER,title_size=14.5,body_size=11.2)
        threats=[
            ("LLM-only judges","Agreement ≠ human correctness",RED),
            ("Same-family confound","Gemini judge cùng họ với 2 tutor",AMBER),
            ("Principle imbalance","Challenge 8; Explanation 863",PURPLE),
            ("Rubric calibration","Questioning/common criteria khó chấm",TEAL),
        ]
        for i,(a,b,c) in enumerate(threats):
            x=6.88+(i%2)*2.94; y=2.94+(i//2)*1.30
            add_card(s,x,y,2.72,1.08,a,b,c,title_size=13,body_size=10.2,number=i+1)
        add_takeaway(s,"Kết quả hỗ trợ relative discrimination, không phải absolute quality certification hay bằng chứng same-family bias.",RED)
        finish(s,"Tables VI–VII; Section VI")

        # 12. Conclusions and meeting decisions
        s=new_slide(prs); add_title(s,"Kết luận và các điểm cần chốt tối nay","Kết luận")
        add_card(s,.72,1.48,3.78,4.70,"Kết luận khoa học",[
            "Pipeline truy vết từ 1.050 hội thoại đến 1.400 mẫu.",
            "6 nguyên tắc + 6 năng lực tạo rubric theo ngữ cảnh.",
            "Phân biệt ổn định khác biệt lớn, chưa phân giải chắc khác biệt nhỏ."],GREEN,title_size=17,body_size=13)
        add_card(s,4.78,1.48,3.78,4.70,"Validation tiếp theo",[
            "HNMU chấm độc lập response pair bằng 22 rubric.",
            "Phân xử thành human reference.",
            "Đo lỗi Gemini/GPT; ưu tiên Questioning, common criteria và principle hiếm."],TEAL,title_size=17,body_size=13)
        add_card(s,8.84,1.48,3.78,4.70,"Cần giáo sư chốt",[
            "Novelty claim và ba đóng góp.",
            "Cách diễn giải 1.400 mẫu eligible.",
            "Reporting hai judge + limitations.",
            "Tác giả, affiliation, corresponding author, acknowledgement."],RED,title_size=17,body_size=13)
        add_takeaway(s,"Nếu ba cột trên được thống nhất, paper sẵn sàng cho vòng sửa cuối và nộp KSE.",NAVY)
        finish(s)

        prs.save(OUT_PPTX)
        print(f"Wrote {OUT_PPTX} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build_deck()
